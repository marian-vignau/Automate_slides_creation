#!/usr/bin/fades
import warnings
import random
import json
from pathlib import Path
import sys

from pptx import Presentation  # fades python-pptx
from pptx.util import Inches
from process_template import Layouts, Layout
from rich.console import Console  # fades rich

import logging

console = Console()

# --- Set up warning logging to file ---
log_file = "warnings.log"

# Configure the logging system
logging.basicConfig(
    level=logging.WARNING,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        logging.FileHandler(log_file),
    ],
)
# Route warnings through the logging system
logging.captureWarnings(True)


def _null_warning(*args, **kwargs): ...


warnings.showwarning = _null_warning  # Suppress all direct warning outputs


def output(*args):
    for arg in args:
        console.print(arg)


def add_content(plh, content):
    text_frame = plh.text_frame
    text_frame.clear()
    for line in content:
        line = line.strip()
        p = text_frame.add_paragraph()
        if line.startswith("-"):
            text = line[1:].strip()  # Remove '- '
            p.level = 0
        else:
            text = line
        bold = False
        for clean_text in text.split("**"):
            # Basic markdown bold handling
            text = clean_text.replace("*", " ").strip()
            if text:
                run = p.add_run()
                run.text = text + " "
                run.font.bold = bold
            bold = not bold


def add_images_notes(slide, slide_data):
    notes = slide_data.get("notes", [])
    visuals = slide_data.get("visual", [])
    if visuals and len(slide.shapes) > 0:
        # Add image to the right of content (example placement)
        left = Inches(5)
        top = Inches(1.5)
        try:
            slide.shapes.add_picture(visuals[0], left, top, height=Inches(3))
        except Exception as e:
            notes.extend(["Visuals:"] + visuals)
            logging.debug(f"Error adding image: {e}")
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "\n".join(notes)


def add_images_notes_(slide, slide_data):
    notes = []
    for key, value in slide_data.items():
        notes.append(f"{key}:")
        notes.extend(value)
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "\n".join(notes)


def on_slide(part, value):
    if len(value) and part not in ["notes", "visual", "index"]:
        return True
    return False


def add_slide(prs, slide_data, map):
    showed_data = {k: v for k, v in slide_data.items() if on_slide(k, v)}
    layouts = map.get_fitted_layouts(showed_data)
    random.shuffle(layouts)
    layout = layouts[0]
    slide = prs.slides.add_slide(layout.obj)

    layout = Layout.new_from_slide(slide)
    # output("Showed data", showed_data)
    _, placeholders = layout.get_fitting(showed_data)

    # output("placeholders", placeholders)
    for part, v in slide_data.items():
        if not on_slide(part, v):
            continue

        if placeholders.get(part):
            add_content(placeholders[part].obj, slide_data[part])
        else:
            output(f"Placeholder not found: {part} for slide {slide_data['index']}")
            slide_data["notes"].extend(
                ["Placeholder not found:"] + [part] + slide_data[part]
            )
    if slide:
        add_images_notes(slide, slide_data)
        return 1
    return 0


def create_presentation(template_path, json_path, output_path):
    # Load the template presentation
    prs = Presentation(template_path)
    map = Layouts(template_path)

    previous_slides = len(prs.slides)
    if previous_slides > 0:
        output(
            f"Warning: Template already contains {previous_slides} slides.",
            "Cannot be removed.",
        )

    # Load the JSON data
    with open(json_path, "r") as f:
        slides_data = json.load(f)

    new_slides = 0
    for idx, slide_data in enumerate(slides_data):
        logging.debug(
            f"Slide {idx + 1}/{len(slides_data)} {' '.join(slide_data.keys())}"
        )
        new_slides += add_slide(prs, slide_data, map)

    # Save the presentation
    prs.save(output_path)
    return dict(
        loaded_data=len(slides_data),
        new_slides=new_slides,
        total_slides=len(prs.slides),
        previous_slides=previous_slides,
    )


# Usage
if __name__ == "__main__":
    if len(sys.argv) != 4:
        output(
            "Usage: python create_w_template.py <source.json> <template.potx> <output.pptx>"
        )
        sys.exit(1)

    source_json = sys.argv[1]
    template_ppt = sys.argv[2]
    output_ppt = sys.argv[3]
    fn = Path(source_json).expanduser().resolve()
    if not fn.exists():
        sys.exit(f"Source JSON file not found: {source_json}")

    fn = Path(template_ppt).expanduser().resolve()
    if not fn.exists():
        sys.exit(f"Template file not found: {template_ppt}")
    stats = create_presentation(
        template_path=str(fn), json_path=source_json, output_path=output_ppt
    )
    output(stats, f"Template applied and saved to: {output_ppt}")
