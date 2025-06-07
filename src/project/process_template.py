#!/usr/bin/fades
import pprint
from collections import defaultdict
import sys
from pathlib import Path
import pptx  # fades python-pptx

replaces = """
only,with,and,section=
_object,_picture=_column
description,text,picture,body,one_column=content
big number,caption,header,main point=title
"""
repla = dict()
for line in replaces.split("\n"):
    if "=" in line:
        value, key = line.split("=")
        repla[key] = value.split(",")


def log(*args, **kwargs): ...


def simplify_name(name):
    name = name.lower()
    for c in "0123456789_-.,":
        name = name.replace(c, " ")
    name = name.strip() + " "

    for card in "one,two,three,four,five,six,seven,eight,nine".split(","):
        if card + " " in name:
            name = name.replace(card + " ", card + "_")
    name = name.replace("s ", " ")
    for key, values in repla.items():
        for value in values:
            if value in name:
                name = name.replace(value, key)

    words = set(x.strip() for x in name.split(" ") if x.strip())
    words = sorted(list(words))
    return tuple(words)


class Layouts:
    def __init__(self, filename):
        self.map = self.process_presentation(filename)

    def process_presentation(self, filename):
        prs = pptx.Presentation(filename)
        map = defaultdict(list)
        for idx, layout in enumerate(prs.slide_masters):
            obj = layout.slide_layouts[0]
            if len(obj.shapes) == 0 and len(obj.placeholders) == 0:
                log(f"{obj.name} dont have shapes")
                continue
            key = simplify_name(obj.name)
            map[key].append((obj, "slide_masters", idx))
        for idx, obj in enumerate(prs.slide_layouts):
            if len(obj.shapes) == 0 and len(obj.placeholders) == 0:
                log(f"{obj.name} dont have shapes")
                continue
            key = simplify_name(obj.name)
            map[key].append((obj, "slide_layouts", idx))
        return map

    def get_layouts(self, parts):
        fit_layouts = []
        n = len(parts)
        for key, values in self.map.items():
            if len(key) == n:
                if n == len(set(parts) & set(key)):
                    fit_layouts.extend(values)
        return [o[0] for o in fit_layouts]

    def __str__(self):
        vmap = {
            " ".join(k): [f"{o[1]}#{o[2]}: {o[0].name}" for o in v]
            for k, v in self.map.items()
        }
        return pprint.pformat(vmap)


if __name__ == "__main__":
    if len(sys.argv) != 2:
        sys.exit("Usage: process_template.py <template.pptx>")

    fn = Path(sys.argv[1])
    fn = fn.expanduser().resolve()
    if not fn.exists():
        sys.exit(f"File not found: {fn}")
    map = Layouts(fn)
    print(map)

    fit = map.get_layouts(["title", "content"])
    print([o.name for o in fit])
