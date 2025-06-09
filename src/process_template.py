#!/usr/bin/fades
import pprint
from dataclasses import dataclass
from collections import defaultdict
import sys
from pathlib import Path
import pptx  # fades python-pptx
from rich.console import Console  # fades rich

console = Console()


def output(*args): ...


# for arg in args:
#     console.print(arg)
#


@dataclass
class Placeholder:
    obj: object
    index: int
    space: int
    typ: str

    @classmethod
    def create(cls, obj, index):
        new = cls(obj, index, 0, "placeholder")
        try:
            new.space = len(obj.text_frame.paragraphs)
        except AttributeError:
            new.space = 0
        new.typ = obj.placeholder_format.type.name.lower()
        if new.typ in ["title", "center_title", "subtitle", "body", "content"]:
            if new.typ == "body":
                new.typ = "content"
            return new
        return None

    @staticmethod
    def type_match(orig, new):
        content = ["subtitle", "body", "content"]
        if orig == new:
            return 100
        if orig in content and new in content:
            if orig == content[0] or new == content[0]:
                return 70
            return 80
        if "title" in orig and "title" in new:
            return 80
        return 0

    def __repr__(self):
        return f"<PH> {self.typ}#{self.index} {self.obj.name}"


@dataclass
class Layout:
    name: str
    typ: str
    position: int
    obj: object

    def __repr__(self):
        s = f"<Layout> {self.name} {self.typ}#{self.position}\n"
        return s + pprint.pformat(self.placeholders)

    @property
    def placeholders(self):
        phds = self.obj.placeholders
        places = {}
        for idx, ph in enumerate(phds):
            pho = Placeholder.create(ph, idx)
            if pho is not None:
                places[idx] = pho
        return places

    @classmethod
    def new_layout(cls, typ, position, obj):
        name = obj.name.strip().lower()
        return Layout(name=name, typ=typ, position=position, obj=obj)

    @classmethod
    def new_from_slide(cls, slide):
        name = slide.name.strip().lower()
        return Layout(name=name, typ="slide", position=0, obj=slide)

    def get_fitting(self, data):
        fit = 0
        fitted = {}
        matched_ph = {}

        for orig in data.keys():
            candidates = {}
            pho = None
            for idx, pho in self.placeholders.items():
                if idx in fitted:
                    continue
                candidates[idx] = Placeholder.type_match(orig, pho.typ)
            if candidates:
                max_candidate = max(candidates.values())
                for idx, fit in candidates.items():
                    if fit == max_candidate:
                        matched_ph[orig] = self.placeholders[idx]
                        fitted[idx] = max_candidate
                        break

        for orig in data.keys():
            if orig not in matched_ph:
                return 0, matched_ph

        return sum(fitted.values()), matched_ph


class Layouts:
    def __init__(self, filename):
        self.map = self.process_presentation(filename)

    def process_presentation(self, filename):
        prs = pptx.Presentation(filename)
        map = []
        for idx, layout in enumerate(prs.slide_masters):
            obj = layout.slide_layouts[0]
            layout = Layout.new_layout("slide_master", idx, obj)
            if layout:
                map.append(layout)
        for idx, obj in enumerate(prs.slide_layouts):
            layout = Layout.new_layout("slide_layout", idx, obj)
            if layout:
                map.append(layout)
        return map

    def get_fitted_layouts(self, slide_data):
        fitting_map = defaultdict(list)
        for layout in self.map:
            fit, _ = layout.get_fitting(slide_data)
            fitting_map[fit].append(layout)
            output(f"=== {layout.name} {fit}", _)
            output(layout)
        max_fit = max(fitting_map.keys())
        return fitting_map[max_fit]

    def __str__(self):
        return pprint.pformat([str(i) for i in self.map])


if __name__ == "__main__":
    if len(sys.argv) != 2:
        sys.exit("Usage: process_template.py <template.pptx>")

    fn = Path(sys.argv[1])
    fn = fn.expanduser().resolve()
    if not fn.exists():
        sys.exit(f"File not found: {fn}")
    map = Layouts(fn)

    slide_data = dict(
        title=["title of the slide"], content=["content line", "content tow"]
    )
    fit = map.get_layouts(slide_data)
    pprint.pprint([str(i) for i in fit])

    pprint.pprint("l fin")
