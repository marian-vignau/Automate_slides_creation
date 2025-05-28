#!/usr/bin/fades
import random
import json
from pathlib import Path
import sys

from pptx import Presentation  # fades python-pptx
from pptx.util import Inches
from process_template import Layouts


def log(*args, **kwargs): ...


def add_title(slide, slide_data):
    title_parts = slide_data.get("title", [])
    if not title_parts:
        return

    title_shape = slide.shapes.title
    title_shape.text = title_parts[0].lstrip("# ").strip()

    if len(title_parts) > 1 and len(slide.placeholders) > 1:
        try:
            for subtitle_placeholder in slide.placeholders:
                ...

            # subtitle_placeholder = slide.placeholders[-1]
            subtitle_placeholder.text = title_parts[1]
        except IndexError:
            print("No subtitle placeholder found")


def search_placeholders(obj):
    subtitle_placeholder = None
    for subtitle_placeholder in obj.placeholders:
        ...
    return subtitle_placeholder


def _search_placeholders(obj):
    content_placeholder = None
    for shape in obj.shapes:
        if shape.has_text_frame and shape.is_placeholder:
            # if shape.placeholder_format.idx == 1:  # Content placeholder
            content_placeholder = shape
    return content_placeholder


def create_presentation(template_path, json_path, output_path):
    # Load the template presentation
    prs = Presentation(template_path)
    map = Layouts(template_path)
    log(map)

    if len(prs.slides) > 0:
        print("Warning: Template already contains slides. I can't remove slides...")

    # Load the JSON data
    with open(json_path, "r") as f:
        slides_data = json.load(f)

    for slide_data in slides_data:
        # Determine slide layout
        # Choose layout based on title (assuming first slide is title slide)
        title_parts = slide_data.get("title", [])
        if len(title_parts) > 1 and "Title Slide" in title_parts[1]:
            layout_type = "title"
        else:
            layout_type = "title content"
        layouts = map.get_layouts(layout_type.split(" "))

        if layout_type == "title":
            layout = layouts[0]
            slide = prs.slides.add_slide(layout)
            add_title(slide, slide_data)
        # Add content (bullet points)
        content = slide_data.get("content", [])
        if content:
            random.shuffle(layouts)
            for layout in layouts:
                plh = search_placeholders(layout)
                if plh:
                    break

            if not plh:
                print(
                    f"No placeholder in {layout.name} found for content: {slide_data}"
                )

            else:
                slide = prs.slides.add_slide(layout)
                # Set title and subtitle
                add_title(slide, slide_data)

                plh = search_placeholders(slide)
                text_frame = plh.text_frame
                text_frame.clear()
                for line in content:
                    line = line.strip()
                    p = text_frame.add_paragraph()
                    if line.startswith("-"):
                        text = line[1:].strip()  # Remove '- '
                        p.level = 0
                    else:
                        text = line
                    bold = False
                    for clean_text in text.split("**"):
                        # Basic markdown bold handling
                        run = p.add_run()
                        run.text = clean_text.replace("*", " ").strip() + " "
                        run.font.bold = bold
                        bold = not bold

        notes = slide_data.get("notes", [])
        # Add images (if any)
        visuals = slide_data.get("visual", [])
        if visuals and len(slide.shapes) > 0:
            # Add image to the right of content (example placement)
            left = Inches(5)
            top = Inches(1.5)
            try:
                slide.shapes.add_picture(visuals[0], left, top, height=Inches(3))
            except Exception as e:
                notes.extend(["Visuals:"] + visuals)
                print(f"Error adding image: {e}")

        # Add speaker notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "\n".join(notes)

    # Save the presentation
    prs.save(output_path)


# Usage
if __name__ == "__main__":
    if len(sys.argv) != 4:
        print(
            "Usage: python create_w_template.py <source.json> <template.potx> <output.pptx>"
        )
        sys.exit(1)

    source_json = sys.argv[1]
    template_ppt = sys.argv[2]
    output_ppt = sys.argv[3]
    fn = Path(template_ppt).expanduser().resolve()
    if not fn.exists():
        sys.exit(f"Template file not found: {template_ppt}")

    create_presentation(
        template_path=str(fn), json_path=source_json, output_path=output_ppt
    )
    print(f"Template applied and saved to: {output_ppt}")
